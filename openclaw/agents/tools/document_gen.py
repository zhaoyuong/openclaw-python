"""Document generation tools for PPT and PDF"""

import logging
import json
from pathlib import Path
from typing import Any
from datetime import datetime

from .base import AgentTool, ToolResult

logger = logging.getLogger(__name__)


class PPTGeneratorTool(AgentTool):
    """Generate PowerPoint presentations"""

    def __init__(self):
        super().__init__()
        self.name = "ppt_generate"
        self.description = (
            "Generate PowerPoint presentations (.pptx files). "
            "Can create slides with titles, content, bullet points, and images. "
            "Perfect for creating presentations, reports, and visual documents. "
            "Files are saved and can be sent via Telegram or other channels."
        )

    def get_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "Presentation title"
                },
                "slides": {
                    "type": "array",
                    "description": "List of slides",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "Slide title"},
                            "content": {
                                "type": "string",
                                "description": "Slide content (supports bullet points with - prefix)"
                            },
                            "layout": {
                                "type": "string",
                                "enum": ["title", "content", "two_column", "blank"],
                                "description": "Slide layout type"
                            }
                        }
                    }
                },
                "author": {
                    "type": "string",
                    "description": "Presentation author (optional)"
                },
                "filename": {
                    "type": "string",
                    "description": "Output filename (optional, auto-generated if not provided)"
                }
            },
            "required": ["title", "slides"]
        }

    async def execute(self, params: dict[str, Any]) -> ToolResult:
        """Generate PowerPoint presentation"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return ToolResult(
                success=False,
                content="",
                error="python-pptx not installed. Run: pip install python-pptx"
            )

        title = params.get("title", "Untitled Presentation")
        slides_data = params.get("slides", [])
        author = params.get("author", "OpenClaw")
        filename = params.get("filename")

        if not slides_data:
            return ToolResult(
                success=False,
                content="",
                error="At least one slide is required"
            )

        try:
            # Create presentation
            prs = Presentation()
            prs.core_properties.title = title
            prs.core_properties.author = author

            # Add slides
            for i, slide_data in enumerate(slides_data):
                slide_title = slide_data.get("title", f"Slide {i+1}")
                slide_content = slide_data.get("content", "")
                layout_type = slide_data.get("layout", "content" if i > 0 else "title")

                if layout_type == "title" or i == 0:
                    # Title slide
                    slide = prs.slides.add_slide(prs.slide_layouts[0])
                    slide.shapes.title.text = slide_title
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = slide_content
                else:
                    # Content slide
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = slide_title
                    
                    if slide_content:
                        body_shape = slide.shapes.placeholders[1]
                        text_frame = body_shape.text_frame
                        text_frame.clear()
                        
                        # Parse bullet points
                        lines = slide_content.split('\n')
                        for line in lines:
                            line = line.strip()
                            if line:
                                # Remove leading bullet markers
                                line = line.lstrip('•-*').strip()
                                p = text_frame.add_paragraph()
                                p.text = line
                                p.level = 0

            # Generate filename
            if not filename:
                safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '_')).strip()
                safe_title = safe_title.replace(' ', '_')[:50]
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"{safe_title}_{timestamp}.pptx"
            
            if not filename.endswith('.pptx'):
                filename += '.pptx'

            # Save to workspace
            output_path = Path.home() / ".openclaw" / "workspace" / "presentations"
            output_path.mkdir(parents=True, exist_ok=True)
            
            file_path = output_path / filename
            prs.save(str(file_path))

            logger.info(f"Generated presentation: {file_path}")

            return ToolResult(
                success=True,
                content=f"✅ Created presentation: {filename}\n\nPath: {file_path}\n\nSlides: {len(slides_data)}",
                metadata={
                    "path": str(file_path),
                    "filename": filename,
                    "slides": len(slides_data),
                    "title": title
                }
            )

        except Exception as e:
            logger.error(f"PPT generation failed: {e}", exc_info=True)
            return ToolResult(
                success=False,
                content="",
                error=f"Failed to generate presentation: {str(e)}"
            )


class PDFGeneratorTool(AgentTool):
    """Generate PDF documents"""

    def __init__(self):
        super().__init__()
        self.name = "pdf_generate"
        self.description = (
            "Generate PDF documents from text content or markdown. "
            "Can create reports, documents, and formatted PDFs. "
            "Supports basic formatting and is useful for creating printable documents."
        )

    def get_schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "content": {
                    "type": "string",
                    "description": "Content to convert to PDF (supports markdown)"
                },
                "title": {
                    "type": "string",
                    "description": "Document title (optional)"
                },
                "filename": {
                    "type": "string",
                    "description": "Output filename (optional, auto-generated if not provided)"
                }
            },
            "required": ["content"]
        }

    async def execute(self, params: dict[str, Any]) -> ToolResult:
        """Generate PDF document"""
        content = params.get("content", "")
        title = params.get("title", "Document")
        filename = params.get("filename")

        if not content:
            return ToolResult(
                success=False,
                content="",
                error="Content is required"
            )

        try:
            # Try using reportlab for PDF generation
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                from reportlab.lib.units import inch
            except ImportError:
                return ToolResult(
                    success=False,
                    content="",
                    error="reportlab not installed. Run: pip install reportlab"
                )

            # Generate filename
            if not filename:
                safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '_')).strip()
                safe_title = safe_title.replace(' ', '_')[:50]
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"{safe_title}_{timestamp}.pdf"
            
            if not filename.endswith('.pdf'):
                filename += '.pdf'

            # Save to workspace
            output_path = Path.home() / ".openclaw" / "workspace" / "documents"
            output_path.mkdir(parents=True, exist_ok=True)
            
            file_path = output_path / filename

            # Create PDF
            doc = SimpleDocTemplate(str(file_path), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            # Add title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.2 * inch))

            # Add content (split by paragraphs)
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    p = Paragraph(para.replace('\n', '<br/>'), styles['BodyText'])
                    story.append(p)
                    story.append(Spacer(1, 0.1 * inch))

            # Build PDF
            doc.build(story)

            logger.info(f"Generated PDF: {file_path}")

            return ToolResult(
                success=True,
                content=f"✅ Created PDF: {filename}\n\nPath: {file_path}",
                metadata={
                    "path": str(file_path),
                    "filename": filename,
                    "title": title
                }
            )

        except Exception as e:
            logger.error(f"PDF generation failed: {e}", exc_info=True)
            return ToolResult(
                success=False,
                content="",
                error=f"Failed to generate PDF: {str(e)}"
            )
