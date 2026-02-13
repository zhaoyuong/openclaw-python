#!/usr/bin/env python3
"""
PowerPoint Generation Script

Generates .pptx files from JSON configuration.

Usage:
    generate_ppt.py --config <config.json> --output <output.pptx> [--template <template.pptx>]

Configuration JSON format:
{
    "title": "Presentation Title",
    "author": "Author Name",
    "slides": [
        {
            "layout": "title",
            "title": "Main Title",
            "subtitle": "Subtitle"
        },
        {
            "layout": "content",
            "title": "Slide Title",
            "content": {
                "bullets": ["Point 1", "Point 2"],
                "text": "Alternative text content"
            }
        }
    ]
}
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx>=0.6.23", file=sys.stderr)
    sys.exit(1)


def create_title_slide(prs: Presentation, slide_def: dict[str, Any]) -> None:
    """Create a title slide (layout 0)"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = slide_def.get("title", "")
    subtitle.text = slide_def.get("subtitle", "")


def create_content_slide(prs: Presentation, slide_def: dict[str, Any]) -> None:
    """Create a content slide with title and content (layout 1)"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = slide_def.get("title", "")
    
    content = slide_def.get("content", {})
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Handle bullet points
    if "bullets" in content:
        for i, bullet_text in enumerate(content["bullets"]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
    
    # Handle plain text
    elif "text" in content:
        p = text_frame.paragraphs[0]
        p.text = content["text"]


def create_section_slide(prs: Presentation, slide_def: dict[str, Any]) -> None:
    """Create a section header slide (layout 2)"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    
    title = slide.shapes.title
    title.text = slide_def.get("title", "")


def create_blank_slide(prs: Presentation, slide_def: dict[str, Any]) -> None:
    """Create a blank slide (layout 6)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title if provided
    if "title" in slide_def:
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = slide_def["title"]
        p.font.size = Pt(32)
        p.font.bold = True


def load_template(template_path: str | None) -> Presentation:
    """Load presentation template"""
    if template_path and Path(template_path).exists():
        return Presentation(template_path)
    
    # Check for default template in assets/
    script_dir = Path(__file__).parent
    default_template = script_dir.parent / "assets" / "default_template.pptx"
    
    if default_template.exists():
        return Presentation(str(default_template))
    
    # Create blank presentation
    return Presentation()


def generate_presentation(config: dict[str, Any], output_path: str, template_path: str | None = None) -> None:
    """Generate PowerPoint presentation from configuration"""
    
    # Load template or create new presentation
    prs = load_template(template_path)
    
    # Set presentation properties
    if "title" in config:
        prs.core_properties.title = config["title"]
    if "author" in config:
        prs.core_properties.author = config["author"]
    
    # Generate slides
    slides = config.get("slides", [])
    for slide_def in slides:
        layout = slide_def.get("layout", "content")
        
        if layout == "title":
            create_title_slide(prs, slide_def)
        elif layout == "content":
            create_content_slide(prs, slide_def)
        elif layout == "section":
            create_section_slide(prs, slide_def)
        elif layout == "blank":
            create_blank_slide(prs, slide_def)
        else:
            print(f"Warning: Unknown layout '{layout}', using content layout", file=sys.stderr)
            create_content_slide(prs, slide_def)
    
    # Save presentation
    output_file = Path(output_path).expanduser()
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    
    return str(output_file.absolute())


def main():
    parser = argparse.ArgumentParser(
        description="Generate PowerPoint presentations from JSON configuration",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Example:
    generate_ppt.py --config slides.json --output presentation.pptx

Configuration format:
    {
        "title": "My Presentation",
        "author": "John Doe",
        "slides": [
            {"layout": "title", "title": "Main Title", "subtitle": "Subtitle"},
            {"layout": "content", "title": "Content", "content": {"bullets": ["A", "B"]}}
        ]
    }
        """
    )
    
    parser.add_argument(
        "--config",
        required=True,
        help="Path to JSON configuration file"
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Output .pptx file path"
    )
    parser.add_argument(
        "--template",
        help="Optional template .pptx file"
    )
    parser.add_argument(
        "--workspace",
        help="Session workspace directory (for relative output paths)"
    )
    
    args = parser.parse_args()
    
    try:
        # Load configuration
        config_path = Path(args.config).expanduser()
        if not config_path.exists():
            print(f"Error: Configuration file not found: {config_path}", file=sys.stderr)
            sys.exit(1)
        
        with open(config_path, "r", encoding="utf-8") as f:
            config = json.load(f)
        
        # Resolve output path
        # If workspace is provided and output is relative, save to workspace
        output_path = Path(args.output)
        
        if not output_path.is_absolute() and args.workspace:
            # Relative path + workspace -> save to session workspace
            workspace = Path(args.workspace).expanduser()
            output_file_path = workspace / output_path
        else:
            # Absolute path -> use as-is
            output_file_path = output_path.expanduser()
        
        # Generate presentation
        output_file = generate_presentation(config, str(output_file_path), args.template)
        
        # Print JSON output for agent to capture (triggers auto file sending)
        result = {
            "file_path": str(output_file),
            "file_type": "document",
            "caption": f"Generated presentation: {config.get('title', 'Presentation')}"
        }
        print(json.dumps(result, ensure_ascii=False))
        sys.exit(0)
        
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON configuration: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
