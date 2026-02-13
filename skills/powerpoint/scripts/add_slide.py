#!/usr/bin/env python3
"""
Add Slide Helper Script

Adds a single slide to an existing PowerPoint presentation.

Usage:
    add_slide.py --pptx <presentation.pptx> --layout <type> --title <text> [--content <json>] --output <output.pptx>
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx>=0.6.23", file=sys.stderr)
    sys.exit(1)


def add_slide_to_presentation(
    pptx_path: str,
    layout: str,
    title: str,
    content: dict | None = None,
    output_path: str | None = None
) -> str:
    """Add a slide to an existing presentation"""
    
    # Load presentation
    prs_path = Path(pptx_path).expanduser()
    if not prs_path.exists():
        raise FileNotFoundError(f"Presentation file not found: {prs_path}")
    
    prs = Presentation(str(prs_path))
    
    # Determine layout index
    layout_map = {
        "title": 0,
        "content": 1,
        "section": 2,
        "blank": 6
    }
    layout_idx = layout_map.get(layout, 1)
    
    # Add slide
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set title
    if hasattr(slide, "shapes") and hasattr(slide.shapes, "title"):
        slide.shapes.title.text = title
    
    # Add content if provided
    if content and layout == "content":
        if 1 < len(slide.placeholders):
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()
            
            # Handle bullets
            if "bullets" in content:
                for i, bullet_text in enumerate(content["bullets"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet_text
                    p.level = 0
            
            # Handle text
            elif "text" in content:
                p = text_frame.paragraphs[0]
                p.text = content["text"]
    
    # Save presentation
    if output_path is None:
        output_path = pptx_path
    
    output_file = Path(output_path).expanduser()
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    
    return str(output_file.absolute())


def main():
    parser = argparse.ArgumentParser(
        description="Add a slide to an existing PowerPoint presentation"
    )
    
    parser.add_argument(
        "--pptx",
        required=True,
        help="Path to existing presentation file"
    )
    parser.add_argument(
        "--layout",
        required=True,
        choices=["title", "content", "section", "blank"],
        help="Slide layout type"
    )
    parser.add_argument(
        "--title",
        required=True,
        help="Slide title"
    )
    parser.add_argument(
        "--content",
        help="Slide content as JSON (for content layout)"
    )
    parser.add_argument(
        "--output",
        help="Output file path (defaults to input file)"
    )
    
    args = parser.parse_args()
    
    try:
        # Parse content if provided
        content = None
        if args.content:
            content = json.loads(args.content)
        
        # Add slide
        output_file = add_slide_to_presentation(
            args.pptx,
            args.layout,
            args.title,
            content,
            args.output
        )
        
        print(output_file)
        sys.exit(0)
        
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON content: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
