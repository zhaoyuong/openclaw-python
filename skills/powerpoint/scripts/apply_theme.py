#!/usr/bin/env python3
"""
Apply Theme Helper Script

Applies a theme from a template presentation to an existing presentation.

Usage:
    apply_theme.py --pptx <presentation.pptx> --theme <theme.pptx> --output <output.pptx>

Note: This copies the slide master and layouts from the theme presentation.
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx>=0.6.23", file=sys.stderr)
    sys.exit(1)


def apply_theme(pptx_path: str, theme_path: str, output_path: str) -> str:
    """Apply a theme to a presentation
    
    Note: python-pptx doesn't directly support copying slide masters,
    so this function creates a new presentation from the theme and
    copies slide content over.
    """
    
    # Load source presentation and theme
    prs_path = Path(pptx_path).expanduser()
    theme_file = Path(theme_path).expanduser()
    
    if not prs_path.exists():
        raise FileNotFoundError(f"Presentation file not found: {prs_path}")
    if not theme_file.exists():
        raise FileNotFoundError(f"Theme file not found: {theme_file}")
    
    source_prs = Presentation(str(prs_path))
    themed_prs = Presentation(str(theme_file))
    
    # Copy slides from source to themed presentation
    for source_slide in source_prs.slides:
        # Determine layout (use content layout as default)
        slide_layout = themed_prs.slide_layouts[1]
        
        # Create new slide with themed layout
        new_slide = themed_prs.slides.add_slide(slide_layout)
        
        # Copy shapes from source slide
        for shape in source_slide.shapes:
            if hasattr(shape, "text"):
                # Find corresponding shape in new slide and copy text
                for new_shape in new_slide.shapes:
                    if hasattr(new_shape, "text") and new_shape.shape_type == shape.shape_type:
                        try:
                            new_shape.text = shape.text
                            break
                        except:
                            pass
    
    # Save themed presentation
    output_file = Path(output_path).expanduser()
    output_file.parent.mkdir(parents=True, exist_ok=True)
    themed_prs.save(str(output_file))
    
    return str(output_file.absolute())


def main():
    parser = argparse.ArgumentParser(
        description="Apply a theme to an existing PowerPoint presentation"
    )
    
    parser.add_argument(
        "--pptx",
        required=True,
        help="Path to presentation file to modify"
    )
    parser.add_argument(
        "--theme",
        required=True,
        help="Path to theme template file"
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Output file path"
    )
    
    args = parser.parse_args()
    
    try:
        output_file = apply_theme(args.pptx, args.theme, args.output)
        print(output_file)
        sys.exit(0)
        
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
